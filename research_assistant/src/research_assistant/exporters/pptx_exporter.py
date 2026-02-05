"""PPTX exporter for markdown-ish reports.

Creates a simple presentation:
- Title slide
- One slide per top-level section (## Heading)
"""

from __future__ import annotations

from dataclasses import dataclass
import re
from typing import Iterable


@dataclass(frozen=True)
class PptxExportOptions:
    title: str
    subtitle: str | None = None


def _iter_lines(text: str) -> Iterable[str]:
    for line in (text or "").splitlines():
        yield line.rstrip("\n")


def _strip_md_inline(text: str) -> str:
    # strip **bold** markers, keep text
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    return text


def export_pptx(markdown_text: str, output_path: str, *, options: PptxExportOptions) -> None:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:  # pragma: no cover
        raise RuntimeError(
            "Missing dependency for PPTX export. Install with:\n"
            "  pip install python-pptx\n"
            "Or:\n"
            "  pip install -e '.[export]'"
        ) from e

    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = options.title
    if options.subtitle and slide.placeholders and len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = options.subtitle
        except Exception:
            pass

    current_title: str | None = None
    current_bullets: list[str] = []

    def flush_section() -> None:
        nonlocal current_title, current_bullets
        if not current_title:
            return
        layout = prs.slide_layouts[1]  # Title and Content
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = current_title
        body = s.shapes.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(current_bullets[:10]):  # keep slides readable
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = b
            p.level = 0
        current_title = None
        current_bullets = []

    for raw in _iter_lines(markdown_text):
        line = raw.strip()
        if not line:
            continue

        if line.startswith("## "):
            flush_section()
            current_title = _strip_md_inline(line[3:].strip())
            continue

        if line.startswith("### "):
            # treat subheading as a bullet marker within section
            if current_title:
                current_bullets.append(_strip_md_inline(line[4:].strip()))
            continue

        if line.startswith("- "):
            if current_title:
                current_bullets.append(_strip_md_inline(line[2:].strip()))
            continue

        m = re.match(r"^(\d+)[\.\)]\s+(.*)$", line)
        if m and current_title:
            current_bullets.append(_strip_md_inline(m.group(2).strip()))
            continue

    flush_section()
    prs.save(output_path)

